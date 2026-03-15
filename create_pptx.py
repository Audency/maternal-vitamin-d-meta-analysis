#!/usr/bin/env python3
"""
Generate academic PowerPoint presentation:
Maternal Vitamin D Status and Fetal Growth / Adiposity — SR & Meta-Analysis
White minimalist background, LaTeX-style formulas as text, Lancet style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
RED    = RGBColor(0xB5, 0x00, 0x1C)
ORANGE = RGBColor(0xD4, 0x63, 0x1A)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
DGREY  = RGBColor(0x44, 0x44, 0x44)
LGREY  = RGBColor(0x90, 0x90, 0x90)
BLACK  = RGBColor(0x11, 0x11, 0x11)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BG = RGBColor(0xF5, 0xF6, 0xF8)

FIG_DIR = os.path.expanduser(
    "~/Desktop/AUDENCIO/Producao artigos/Artigo RS Isabel Vitamina D /Resultados /figures_v6/")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper functions ─────────────────────────────────────────────────────────
def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_bottom_line(slide):
    left = Inches(0.6)
    top = Inches(7.05)
    w = Inches(12.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_slide_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = LGREY
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_white_bg(slide)

    # Navy bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0),
                                      Inches(11.7), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2),
                                       Inches(11.7), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = DGREY
    p2.alignment = PP_ALIGN.LEFT

    # Bottom line
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, Inches(7.0), prs.slide_width, Inches(0.08))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = NAVY
    bar2.line.fill.background()

    return slide


def add_content_slide(title, bullets, slide_num, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_bg(slide)
    add_bottom_line(slide)
    add_slide_number(slide, slide_num)

    # Section title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3),
                                      Inches(12.1), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # Thin line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(1.05),
                                   Inches(12.1), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Bullets
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3),
                                       Inches(11.5), Inches(5.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(17)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
        p.level = 0
        if bullet.startswith("   "):
            p.font.size = Pt(15)
            p.font.color.rgb = DGREY

    if note:
        txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.4),
                                           Inches(11.5), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.font.color.rgb = LGREY

    return slide


def add_figure_slide(title, img_path, slide_num, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_bg(slide)
    add_bottom_line(slide)
    add_slide_number(slide, slide_num)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3),
                                      Inches(12.1), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.95),
                                   Inches(12.1), Pt(1.2))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Image
    if os.path.exists(img_path):
        img = slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.2),
                                        width=Inches(12.1))
        # if image is too tall, constrain
        if img.height > Inches(5.3):
            ratio = Inches(5.3) / img.height
            img.height = Inches(5.3)
            img.width = int(img.width * ratio)
            # center horizontally
            img.left = int((prs.slide_width - img.width) / 2)

    if caption:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(6.6),
                                           Inches(12.1), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.color.rgb = LGREY

    return slide


# =============================================================================
# SLIDES
# =============================================================================

# ── 1. TITLE SLIDE ───────────────────────────────────────────────────────────
add_title_slide(
    "Association Between Maternal Vitamin D Levels\n"
    "and Fetal Growth and Adiposity",
    "A Systematic Review and Meta-Analysis\n\n"
    "Isabel O. A. & Audencio Victor\n"
    "PROSPERO: CRD42024590338"
)

# ── 2. BACKGROUND ───────────────────────────────────────────────────────────
add_content_slide("Background", [
    "\u2022 Vitamin D [25(OH)D] regulates calcium homeostasis, immune function, cell growth, and gene transcription",
    "\u2022 Maternal 25(OH)D crosses the placenta and is essential for fetal skeletal development and growth",
    "\u2022 Deficiency (< 50 nmol/L) affects up to 40\u201360% of pregnant women worldwide",
    "\u2022 Fetal growth outcomes: SGA, FGR, biometry (FL, AC, BPD, HC), estimated fetal weight (EFW), adiposity",
    "\u2022 Existing evidence is inconsistent due to heterogeneity in study design, VitD cutoffs, and analytical methods",
    "\u2022 No comprehensive meta-analysis pooling both observational and RCT data on fetal growth and adiposity",
], 2)

# ── 3. OBJECTIVE ─────────────────────────────────────────────────────────────
add_content_slide("Objective", [
    "\u2022 To synthesize the available evidence on the association between maternal",
    "  25(OH)D concentrations during pregnancy and fetal growth and adiposity",
    "",
    "\u2022 Primary outcome: SGA / FGR (dichotomous)",
    "\u2022 Secondary outcomes:",
    "   \u2013 Fetal biometry / EFW (dichotomous, OR)",
    "   \u2013 Fetal bone measurements (continuous, MD) from RCTs",
    "   \u2013 Fetal adiposity",
], 3)

# ── 4. METHODS — SEARCH ─────────────────────────────────────────────────────
add_content_slide("Methods \u2014 Search Strategy & Eligibility", [
    "\u2022 PRISMA 2020 guidelines; registered at PROSPERO (CRD42024590338)",
    "\u2022 Databases: PubMed, EMBASE, Scopus, LILACS, SciELO, Web of Science",
    "\u2022 Search date: up to October 2024",
    "\u2022 Inclusion: observational studies and RCTs in pregnant women with 25(OH)D measurement",
    "   and fetal growth/biometry/adiposity outcomes",
    "\u2022 Exclusion: multiple pregnancies, reviews, editorials, conference abstracts",
    "\u2022 PECOS: Population (pregnant women), Exposure (25(OH)D status),",
    "   Comparator (sufficient vs deficient), Outcome (fetal growth), Study (obs + RCT)",
    "\u2022 Selection: two independent reviewers (AV, IOA); disagreements resolved by consensus",
], 4)

# ── 5. METHODS — STATISTICAL ────────────────────────────────────────────────
add_content_slide("Methods \u2014 Statistical Analysis", [
    "\u2022 Dichotomous outcomes: Odds Ratio (OR) with inverse-variance weighting",
    "   ln(OR) = ln(a\u00b7d / b\u00b7c);  SE = \u221a(1/a + 1/b + 1/c + 1/d)",
    "",
    "\u2022 Continuous outcomes (RCTs): Mean Difference (MD)",
    "   MD = X\u0304\u2091 \u2212 X\u0304\u1d04;  SE = \u221a(s\u00b2\u2091/n\u2091 + s\u00b2\u1d04/n\u1d04)",
    "",
    "\u2022 Model selection: adaptive approach",
    "   Fixed-effects (Mantel\u2013Haenszel) when I\u00b2 < 50%",
    "   Random-effects (REML) when I\u00b2 \u2265 50%",
    "",
    "\u2022 Heterogeneity: Cochran\u2019s Q statistic and I\u00b2",
    "   I\u00b2 = [(Q \u2212 df) / Q] \u00d7 100%",
    "",
    "\u2022 Publication bias: Egger\u2019s regression, Begg\u2019s rank test, Duval & Tweedie trim-and-fill",
    "\u2022 Quality: JBI Critical Appraisal + NOS (observational); Cochrane RoB 2 (RCTs)",
    "\u2022 Certainty of evidence: GRADE framework",
], 5, note="Software: R (meta, metafor packages)")

# ── 6. METHODS — FORMULAS ───────────────────────────────────────────────────
add_content_slide("Methods \u2014 Key Formulas", [
    "Pooled OR (fixed-effects, inverse-variance):",
    "   ln(OR\u209a) = \u03a3 w\u1d62 \u00b7 ln(OR\u1d62) / \u03a3 w\u1d62 ,  where  w\u1d62 = 1 / SE\u1d62\u00b2",
    "",
    "Random-effects (DerSimonian\u2013Laird / REML):",
    "   w\u1d62* = 1 / (SE\u1d62\u00b2 + \u03c4\u00b2)",
    "   \u03c4\u00b2 estimated via restricted maximum likelihood (REML)",
    "",
    "Heterogeneity:",
    "   Q = \u03a3 w\u1d62(ln OR\u1d62 \u2212 ln OR\u209a)\u00b2 ~ \u03c7\u00b2(k\u22121)",
    "   I\u00b2 = max(0, (Q \u2212 (k\u22121)) / Q \u00d7 100%)",
    "",
    "Risk Difference (clinical interpretation):",
    "   RD = p\u2080 \u00d7 (OR \u2212 1) / [1 + p\u2080 \u00d7 (OR \u2212 1)]",
    "   NNH = 1 / |RD|",
    "",
    "Egger\u2019s test:  SND\u1d62 = ln(OR\u1d62) / SE\u1d62  regressed on  precision\u1d62 = 1/SE\u1d62",
], 6)

# ── 7. FLOW DIAGRAM ─────────────────────────────────────────────────────────
add_content_slide("Results \u2014 Study Selection", [
    "\u2022 Initial search: 3,135 records across 5 databases",
    "\u2022 After duplicate removal: 3,412 screened (phase 1)",
    "\u2022 Full-text assessed: 30 articles (phase 2)",
    "\u2022 Included: 26 studies (24 observational + 2 RCTs)",
    "\u2022 Total participants: 26,542",
    "\u2022 Countries: 16 (Asia 46%, Europe 31%, America 15%, Africa 4%)",
    "\u2022 Publication period: 2010\u20132025",
    "",
    "\u2022 Meta-analysis (pooled): k = 4 (SGA/FGR) + k = 4 (Biometry/EFW)",
    "\u2022 Narrative synthesis: remaining 18 studies",
], 7)

# ── 8. STUDY CHARACTERISTICS ────────────────────────────────────────────────
add_figure_slide("Results \u2014 Study Characteristics",
                 FIG_DIR + "StudyCharacteristics.png", 8)

# ── 9. QUALITY ───────────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 Methodological Quality (JBI + NOS)",
                 FIG_DIR + "Quality_JBI_NOS.png", 9,
                 "JBI Critical Appraisal (Q1\u2013Q9) and Newcastle-Ottawa Scale for 24 observational studies.")

# ── 10. PRIMARY: SGA/FGR ────────────────────────────────────────────────────
add_figure_slide("Results \u2014 Primary Meta-Analysis: SGA/FGR",
                 FIG_DIR + "Forest_SGA_FGR.png", 10,
                 "Pooled OR = 1.95 (95% CI 1.47\u20132.59); I\u00b2 = 0%; Fixed-effects model; k = 4, N = 9,033.")

# ── 11. SECONDARY: BIOMETRY ─────────────────────────────────────────────────
add_figure_slide("Results \u2014 Secondary Meta-Analysis: Biometry/EFW",
                 FIG_DIR + "Forest_Biometry_EFW.png", 11,
                 "Pooled OR = 1.16 (95% CI 1.09\u20131.23); I\u00b2 = 33%; Fixed-effects model; k = 4, N = 26,542.")

# ── 12. RCT ──────────────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 RCT Continuous Outcomes (Mean Difference)",
                 FIG_DIR + "Forest_RCT_MD.png", 12,
                 "Vafaei 2019 (Low RoB): FL +1.98 mm, HL +1.39 mm. Srilekha 2021 (High RoB): EFW +277 g.")

# ── 13. SENSITIVITY ─────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 Leave-One-Out Sensitivity Analysis",
                 FIG_DIR + "LOO_Sensitivity.png", 13,
                 "All LOO estimates remain statistically significant, confirming robustness.")

# ── 14. PUBLICATION BIAS ─────────────────────────────────────────────────────
add_figure_slide("Results \u2014 Publication Bias (Funnel Plot & Trim-and-Fill)",
                 FIG_DIR + "Funnel_TrimFill.png", 14,
                 "Duval & Tweedie trim-and-fill analysis. Egger/Begg tests underpowered (k = 4).")

# ── 15. SUBGROUP ─────────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 Subgroup by Vitamin D Cutoff",
                 FIG_DIR + "Cutoff_Subgroup.png", 15,
                 "All subgroups directionally consistent regardless of deficiency threshold.")

# ── 16. GRADE ────────────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 GRADE Evidence Profile",
                 FIG_DIR + "GRADE_Profile.png", 16,
                 "Certainty: High (Biometry/EFW), Moderate (SGA/FGR, FL/HL, Infant wt.), Low (EFW/BPD RCT, Adiposity).")

# ── 17. BUBBLE ───────────────────────────────────────────────────────────────
add_figure_slide("Results \u2014 25(OH)D Concentration vs Sample Size",
                 FIG_DIR + "BubblePlot.png", 17,
                 "Studies below 50 nmol/L threshold more frequently reported significant associations.")

# ── 18. KEY FINDINGS ─────────────────────────────────────────────────────────
add_content_slide("Summary of Key Findings", [
    "\u2022 Maternal vitamin D deficiency is significantly associated with increased risk of SGA/FGR",
    "   OR = 1.95 (1.47\u20132.59); NNH = 14 (baseline SGA 8%)",
    "",
    "\u2022 Deficiency also associated with altered fetal biometry/EFW",
    "   OR = 1.16 (1.09\u20131.23); GRADE: High certainty",
    "",
    "\u2022 RCT evidence (limited): supplementation improves fetal bone length",
    "   FL: +1.98 mm (p < 0.001); HL: +1.39 mm (p < 0.001)",
    "",
    "\u2022 Results robust across sensitivity analyses (LOO, subgroup, trim-and-fill)",
    "",
    "\u2022 No significant association with fetal adiposity (k = 2; GRADE: Low)",
    "",
    "\u2022 15/26 studies (58%) reported significant associations with fetal growth",
], 18)

# ── 19. STRENGTHS & LIMITATIONS ──────────────────────────────────────────────
add_content_slide("Strengths & Limitations", [
    "Strengths:",
    "   \u2013 Comprehensive search across 6 databases (3,135 records)",
    "   \u2013 Both observational and RCT evidence synthesized",
    "   \u2013 Adaptive model selection based on heterogeneity (I\u00b2)",
    "   \u2013 Multiple sensitivity analyses and GRADE assessment",
    "   \u2013 Dual quality tools (JBI + NOS for observational; RoB 2 for RCTs)",
    "",
    "Limitations:",
    "   \u2013 Small number of poolable studies (k = 4 per analysis)",
    "   \u2013 Heterogeneity in VitD cutoffs and timing of measurement",
    "   \u2013 Only 2 RCTs available (1 with high risk of bias)",
    "   \u2013 Limited evidence on adiposity outcomes",
    "   \u2013 Residual confounding in observational studies",
], 19)

# ── 20. CONCLUSION ───────────────────────────────────────────────────────────
add_content_slide("Conclusion", [
    "\u2022 Maternal vitamin D deficiency during pregnancy is consistently associated",
    "  with impaired fetal growth, particularly increased SGA/FGR risk (OR \u2248 2.0)",
    "",
    "\u2022 Vitamin D supplementation in RCTs shows promising effects on fetal bone growth,",
    "  though evidence remains limited (k = 2)",
    "",
    "\u2022 Clinical implication: screening and correcting vitamin D deficiency in pregnancy",
    "  may contribute to improved fetal growth outcomes",
    "",
    "\u2022 Future research: large-scale RCTs with standardized protocols, consistent cutoffs,",
    "  and adiposity as a primary endpoint are needed",
], 20)

# ── 21. THANK YOU ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5),
                                  Inches(11.7), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0),
                                   Inches(11.7), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "audenciovictor@gmail.com\nhttps://github.com/Audency/maternal-vitamin-d-meta-analysis"
p2.font.size = Pt(16)
p2.font.color.rgb = DGREY
p2.alignment = PP_ALIGN.CENTER

bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               0, Inches(7.0), prs.slide_width, Inches(0.08))
bar2.fill.solid()
bar2.fill.fore_color.rgb = NAVY
bar2.line.fill.background()


# ── SAVE ─────────────────────────────────────────────────────────────────────
out_path = os.path.expanduser(
    "~/Desktop/AUDENCIO/Producao artigos/Artigo RS Isabel Vitamina D /Presentation_VitaminD_SR.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
